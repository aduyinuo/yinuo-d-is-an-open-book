#!/usr/bin/env python3
"""
Build an annotated-slides page from a compiled PDF instead of a .pptx.

    python internal/slides/pdf_to_deck.py slides-source/asu-brown-bag-latex/asu-brown-bag.pdf \
        --notes slides-source/asu-brown-bag.pptx --slug asu-brown-bag

WHY THIS EXISTS

pptx_to_deck.py rebuilds each slide in HTML from the PowerPoint shapes.
That works for text and pictures, but PowerPoint also draws diagrams
from its own shapes -- boxes, arrows, connectors, groups -- and those
carry no image. They are dropped, the layout collapses around the gap,
and the published page reads as scrambled.

Rendering a compiled PDF sidesteps the whole problem: every page is
already exactly what the slide looks like. The page becomes one image,
so nothing can be dropped or shifted.

The trade is that the slide text is no longer HTML, so it is not
selectable and not indexed. Narration under each slide still is, and it
is the narration that carries the searchable prose.

Narration comes from the .pptx notes, matched to pages by slide title
rather than by position: a LaTeX deck emits extra pages for overlays and
for frames that continue, so page N is not slide N.

    pip install pymupdf python-pptx pillow
"""
import argparse
import html
import json
import os
import re
import shutil
import sys

HERE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.abspath(os.path.join(HERE, "..", ".."))
sys.path.insert(0, HERE)

try:
    import pymupdf
except ImportError:
    sys.exit("Missing PyMuPDF.  pip install pymupdf")

from pptx_to_deck import build_html, write_oembed, write_cover, OUT  # noqa: E402


def notes_by_title(pptx_path):
    """[(normalised title, narration)] for every slide that has notes."""
    try:
        from pptx import Presentation
    except ImportError:
        sys.exit("Missing python-pptx.  pip install python-pptx")

    prs = Presentation(pptx_path)
    out = []
    for slide in prs.slides:
        title = ""
        if slide.shapes.title is not None and slide.shapes.title.has_text_frame:
            title = " ".join(slide.shapes.title.text_frame.text.split())
        note = ""
        if slide.has_notes_slide:
            note = slide.notes_slide.notes_text_frame.text.strip()
        out.append((norm(title), note))
    return out


def norm(t):
    """Loose key: case and punctuation vary between PowerPoint and LaTeX."""
    return re.sub(r"[^a-z0-9 ]+", "", " ".join((t or "").lower().split()))


def page_title(page):
    """The topmost line of a page, which the theme sets as the frame title."""
    text = page.get_text().strip().splitlines()
    return norm(text[0]) if text else ""


def build(pdf_path, slug, notes_path, dpi):
    doc = pymupdf.open(pdf_path)
    media_dir = os.path.join(OUT, "media", slug)
    if os.path.isdir(media_dir):
        shutil.rmtree(media_dir)
    os.makedirs(media_dir, exist_ok=True)

    notes = notes_by_title(notes_path) if notes_path else []
    used = set()
    frame = 0   # index of the source slide the current page belongs to

    frags, narr = {}, {}
    zoom = dpi / 72.0
    W = H = None

    for n, page in enumerate(doc, start=1):
        if W is None:
            W, H = page.rect.width, page.rect.height
        pix = page.get_pixmap(matrix=pymupdf.Matrix(zoom, zoom), alpha=False)
        name = "p%03d.png" % n
        pix.save(os.path.join(media_dir, name))

        frags[n] = ('<div class="slide">'
                    '<div class="shp img" style="left:0;top:0;width:100%%;height:100%%">'
                    '<img src="media/%s/%s" alt="" loading="lazy" '
                    'style="width:100%%;height:100%%;object-fit:contain"></div></div>'
                    % (slug, name))

        # Map pages to slides by counting frame starts rather than by
        # matching titles: the LaTeX titles are rewritten from the
        # PowerPoint ones and often differ, so title matching finds only
        # a fraction of them. A frame that runs past one page is marked
        # "(cont.)" by the theme, so any page WITHOUT that marker begins
        # a new frame, and the generator emits exactly one frame per
        # source slide, in order.
        raw_title = (page.get_text().strip().splitlines() or [""])[0]
        text = ""
        if n > 1 and "(cont.)" not in raw_title:
            if frame < len(notes):
                text = notes[frame][1]
                if text.strip():
                    used.add(frame)
            frame += 1
        parts = [p.strip() for p in re.split(r"^\s*---\s*$", text, flags=re.M) if p.strip()]
        narr[n] = "\n\n".join(parts) if parts else text

    title = slug.replace("-", " ").replace("_", " ")
    ratio = (W / H) if H else 16 / 9
    doc.close()

    page_html = build_html(title, frags, narr, ratio, slug)
    with open(os.path.join(OUT, slug + ".html"), "w", encoding="utf-8", newline="\n") as fh:
        fh.write(page_html)
    write_oembed(slug, title, ratio)
    write_cover(slug, title)

    placed = sum(1 for v in narr.values() if v.strip())
    print("wrote %s" % os.path.join(OUT, slug + ".html"))
    print("  %d pages at %d dpi, %d with narration" % (len(frags), dpi, placed))
    if notes:
        print("  %d of %d slide notes matched by title"
              % (len(used), sum(1 for _, x in notes if x.strip())))


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("pdf")
    ap.add_argument("--slug", required=True)
    ap.add_argument("--notes", help="the .pptx to lift narration from")
    ap.add_argument("--dpi", type=int, default=150)
    args = ap.parse_args()
    if not os.path.exists(args.pdf):
        sys.exit("No such file: %s" % args.pdf)
    build(args.pdf, args.slug, args.notes, args.dpi)


if __name__ == "__main__":
    main()
