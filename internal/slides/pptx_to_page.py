#!/usr/bin/env python3
"""
Render a PowerPoint deck into a GitBook page: one image per slide, with the
speaker notes underneath it. Uses only blocks GitBook renders natively, so it
needs no hosting, no plugin, and no embed.

Usage:  python internal/slides/pptx_to_page.py slides-source/<deck>.pptx <page.md>

Requires LibreOffice (soffice) and poppler's pdftoppm on PATH.
"""
import os, re, sys, glob, shutil, subprocess, tempfile

ROOT = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", ".."))
ASSETS = os.path.join(ROOT, "content", ".gitbook", "assets")


def notes(pptx):
    from pptx import Presentation
    out = []
    for s in Presentation(pptx).slides:
        t = ""
        try:
            if s.has_notes_slide:
                t = (s.notes_slide.notes_text_frame.text or "").strip()
        except Exception:
            pass
        out.append(t)
    return out


def render(pptx, slug, dpi=100, width=1400):
    from PIL import Image
    tmp = tempfile.mkdtemp()
    subprocess.run(["soffice", "--headless", "--norestore", "--convert-to", "pdf",
                    os.path.abspath(pptx), "--outdir", tmp],
                   check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL,
                   timeout=1800)
    pdf = glob.glob(os.path.join(tmp, "*.pdf"))[0]
    subprocess.run(["pdftoppm", "-r", str(dpi), "-png", "-aa", "yes",
                    "-aaVector", "yes", pdf, os.path.join(tmp, "p")],
                   check=True, timeout=1800)
    os.makedirs(ASSETS, exist_ok=True)
    names = []
    for i, f in enumerate(sorted(glob.glob(os.path.join(tmp, "p-*.png"))), start=1):
        im = Image.open(f).convert("RGB")
        if im.width > width:
            im = im.resize((width, round(im.height * width / im.width)), Image.LANCZOS)
        name = "slide-%s-%02d.jpg" % (slug, i)
        im.save(os.path.join(ASSETS, name), "JPEG", quality=80,
                optimize=True, progressive=True)
        names.append(name)
    shutil.rmtree(tmp, ignore_errors=True)
    return names


def relpath_to_assets(page_path):
    d = os.path.dirname(os.path.abspath(page_path))
    return os.path.relpath(ASSETS, d).replace(os.sep, "/")


def build(pptx, page_path, title=None):
    slug = os.path.splitext(os.path.basename(pptx))[0]
    imgs = render(pptx, slug)
    ns = notes(pptx)
    rel = relpath_to_assets(page_path)
    out = []
    for i, img in enumerate(imgs):
        n = ns[i] if i < len(ns) else ""
        out.append('<figure><img src="%s/%s" alt="Slide %d"><figcaption><p>%d</p></figcaption></figure>'
                   % (rel, img, i + 1, i + 1))
        body = "\n".join(p for p in n.split("\n") if p.strip())
        if body:
            out.append(body)
        out.append("")
    return "\n\n".join(out), len(imgs)


if __name__ == "__main__":
    pptx, page = sys.argv[1], sys.argv[2]
    md, n = build(pptx, page)
    open(page, "a", encoding="utf-8", newline="\n").write(md)
    print("%d slides appended to %s" % (n, page))
