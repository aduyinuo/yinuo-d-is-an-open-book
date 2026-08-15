"""
Slide primitives for 20-minute research talks.

Follows internal/slides/TALK-RULES.md: message-carrying titles, very few words,
large sans-serif type, everything the speaker would say lives in the notes.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

D  = RGBColor(0x22, 0x45, 0x2F)
M  = RGBColor(0x4A, 0x7C, 0x59)
L  = RGBColor(0x93, 0xB2, 0x94)
VL = RGBColor(0xE6, 0xEF, 0xE4)
INK= RGBColor(0x1F, 0x24, 0x30)
MUT= RGBColor(0x6B, 0x72, 0x80)
OCH= RGBColor(0xA8, 0x84, 0x3C)
WHT= RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Arial"


class Deck:
    def __init__(self, title, subtitle, note=""):
        self.p = Presentation()
        self.p.slide_width  = Inches(13.333)
        self.p.slide_height = Inches(7.5)
        self._title_slide(title, subtitle, note)

    # ---------------------------------------------------------------- helpers
    def _blank(self):
        return self.p.slides.add_slide(self.p.slide_layouts[6])

    def _notes(self, s, text):
        s.notes_slide.notes_text_frame.text = text.strip()

    def _box(self, s, x, y, w, h, fill=None, line=None):
        from pptx.enum.shapes import MSO_SHAPE
        sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(x), Inches(y), Inches(w), Inches(h))
        sh.shadow.inherit = False
        if fill is None:
            sh.fill.background()
        else:
            sh.fill.solid(); sh.fill.fore_color.rgb = fill
        if line is None:
            sh.line.fill.background()
        else:
            sh.line.color.rgb = line; sh.line.width = Pt(1.25)
        sh.text_frame.text = ""
        return sh

    def _text(self, s, x, y, w, h, runs, size=30, color=INK, bold=False,
              align=PP_ALIGN.LEFT, space=10):
        tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame; tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP
        items = runs if isinstance(runs, (list, tuple)) else [runs]
        for i, t in enumerate(items):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.alignment = align
            para.space_after = Pt(space)
            r = para.add_run(); r.text = t
            r.font.size = Pt(size); r.font.bold = bold
            r.font.color.rgb = color; r.font.name = FONT
        return tb

    def _title(self, s, text, color=D):
        self._text(s, 0.62, 0.45, 12.1, 1.5, text, size=40, color=color, bold=True)

    # ----------------------------------------------------------------- slides
    def _title_slide(self, title, subtitle, note):
        s = self._blank()
        bg = self._box(s, -0.1, -0.1, 13.6, 7.7, fill=D)
        bg.shadow.inherit = False
        self._text(s, 1.0, 2.4, 11.3, 2.2, title, size=52, color=WHT, bold=True)
        self._text(s, 1.0, 4.5, 11.3, 1.2, subtitle, size=24, color=L)
        self._notes(s, note)

    def section(self, label, note=""):
        s = self._blank()
        self._box(s, -0.1, -0.1, 13.6, 7.7, fill=VL)
        self._text(s, 1.0, 3.1, 11.3, 1.6, label, size=46, color=D, bold=True)
        self._notes(s, note)
        return s

    def points(self, title, bullets, note="", size=30):
        """A message title with at most a few short lines under it."""
        s = self._blank()
        self._title(s, title)
        y = 2.2
        for b in bullets:
            self._box(s, 0.62, y + 0.18, 0.16, 0.16, fill=M)
            self._text(s, 1.0, y, 11.7, 1.0, b, size=size)
            y += 1.05
        self._notes(s, note)
        return s

    def statement(self, title, line="", note="", color=D):
        """One sentence, large. For the message slides."""
        s = self._blank()
        self._title(s, title, color=color)
        if line:
            self._text(s, 0.62, 2.6, 12.1, 3.2, line, size=34, color=INK)
        self._notes(s, note)
        return s

    def figure(self, title, caption, note="", placeholder=True, label="figure"):
        """A slide whose body is a picture. Draws a marked frame if absent."""
        s = self._blank()
        self._title(s, title)
        if placeholder:
            self._box(s, 1.6, 2.15, 10.1, 4.1, fill=VL, line=L)
            self._text(s, 1.9, 3.9, 9.5, 0.8, "[ %s ]" % label, size=26,
                       color=OCH, align=PP_ALIGN.CENTER)
        if caption:
            self._text(s, 0.62, 6.35, 12.1, 0.8, caption, size=20, color=MUT)
        self._notes(s, note)
        return s

    def two_col(self, title, left_head, left, right_head, right, note=""):
        s = self._blank()
        self._title(s, title)
        for x, head, items, col in ((0.62, left_head, left, M),
                                    (7.0, right_head, right, OCH)):
            self._text(s, x, 2.15, 5.6, 0.7, head, size=26, color=col, bold=True)
            y = 3.0
            for it in items:
                self._text(s, x, y, 5.6, 1.1, "· " + it, size=24)
                y += 1.0
        self._notes(s, note)
        return s

    def placeholder(self, title, what, note=""):
        """An explicitly unfinished slide."""
        s = self._blank()
        self._title(s, title, color=OCH)
        self._box(s, 0.62, 2.2, 12.1, 3.4, fill=None, line=OCH)
        self._text(s, 1.0, 2.6, 11.4, 2.6,
                   ["[PLACEHOLDER]", what], size=28, color=OCH)
        self._notes(s, note)
        return s

    def save(self, path):
        self.p.save(path)
        return path
