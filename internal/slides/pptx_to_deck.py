#!/usr/bin/env python3
"""
Convert a PowerPoint deck into an annotated-slides HTML page.

Usage:  python internal/slides/pptx_to_deck.py            # convert everything in slides-source/
        python internal/slides/pptx_to_deck.py my.pptx    # convert one file

Input   slides-source/<name>.pptx
Output  docs/slides/<name>.html   (self-contained, no build step, no dependencies)

Narration: whatever you type in the slide's Notes pane becomes the text shown
under that slide. Blank lines make paragraphs; **bold** and _italic_ work.
A line starting with '---' splits the notes so the slide reveals in steps.
"""
import os, sys, re, json, html, io, shutil

from pptx import Presentation
from pptx.util import Emu

HERE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.abspath(os.path.join(HERE, "..", ".."))
SRC  = os.path.join(ROOT, "slides-source")
OUT  = os.path.join(ROOT, "docs", "slides")

# ---------------------------------------------------------------- pptx -> html
def emu_pct(v, total):
    return round(100.0 * float(v) / float(total), 3)

def pt_to_cqw(pt, slide_w_emu):
    """Font size in points -> cqw, so text scales with the slide box."""
    slide_w_pt = float(slide_w_emu) / 12700.0
    return round(100.0 * float(pt) / slide_w_pt, 3)

def color_of(run, default=None):
    try:
        c = run.font.color
        if c and c.type is not None and c.rgb is not None:
            return "#%s" % str(c.rgb)
    except Exception:
        pass
    return default

def para_html(p, slide_w, base_pt):
    bits = []
    for r in p.runs:
        t = html.escape(r.text or "", quote=False)
        if not t:
            continue
        sz = r.font.size.pt if r.font.size else base_pt
        style = ["font-size:%.2fcqw" % pt_to_cqw(sz, slide_w)]
        if r.font.bold:   style.append("font-weight:650")
        if r.font.italic: style.append("font-style:italic")
        col = color_of(r)
        if col: style.append("color:%s" % col)
        bits.append('<span style="%s">%s</span>' % (";".join(style), t))
    if not bits:
        return ""
    align = {1: "center", 2: "right", 3: "justify"}.get(
        getattr(p.alignment, "value", None), "left")
    indent = 2.2 * (p.level or 0)
    return ('<p style="margin:0 0 .45em 0;text-align:%s;margin-left:%.1fcqw">%s</p>'
            % (align, indent, "".join(bits)))

def save_image(sh, media_dir, media_url, counter):
    """Write the picture to its own file, downscaled, and return the URL."""
    try:
        img = sh.image
    except Exception:
        return None
    ext = (img.ext or "png").lower()
    name = "img%03d.%s" % (counter, "jpg" if ext in ("jpeg", "jpg") else ext)
    path = os.path.join(media_dir, name)
    blob = img.blob
    try:
        from PIL import Image
        im = Image.open(io.BytesIO(blob))
        if im.width > 1600:
            im = im.resize((1600, round(im.height * 1600 / im.width)), Image.LANCZOS)
        if im.mode in ("RGBA", "LA", "P") and "transparency" in im.info or im.mode == "RGBA":
            im.save(path if path.endswith(".png") else path + ".png", optimize=True)
            if not path.endswith(".png"):
                name += ".png"; path += ".png"
        else:
            im = im.convert("RGB")
            name = os.path.splitext(name)[0] + ".jpg"
            path = os.path.join(media_dir, name)
            im.save(path, "JPEG", quality=82, optimize=True, progressive=True)
    except Exception:
        with open(path, "wb") as fh:
            fh.write(blob)
    return media_url + "/" + name


def shape_html(sh, slide_w, slide_h, media, idx):
    try:
        left, top = sh.left, sh.top
        width, height = sh.width, sh.height
        if left is None or top is None:
            return ""
    except Exception:
        return ""
    pos = ("left:%.3f%%;top:%.3f%%;width:%.3f%%;height:%.3f%%"
           % (emu_pct(left, slide_w), emu_pct(top, slide_h),
              emu_pct(width, slide_w), emu_pct(height, slide_h)))

    if sh.shape_type == 13 or sh.__class__.__name__ == "Picture":
        url = save_image(sh, media[0], media[1], idx) if media else None
        if not url:
            return ""
        return ('<div class="shp img" style="%s">'
                '<img src="%s" alt="" loading="lazy"></div>' % (pos, url))

    if not sh.has_text_frame:
        # keep simple filled shapes as blocks so layout still reads
        try:
            fill = sh.fill
            if fill.type is not None and fill.fore_color.rgb is not None:
                return ('<div class="shp box" style="%s;background:#%s"></div>'
                        % (pos, str(fill.fore_color.rgb)))
        except Exception:
            pass
        return ""

    tf = sh.text_frame
    if not (tf.text or "").strip():
        return ""
    base_pt = 18
    paras = "".join(para_html(p, slide_w, base_pt) for p in tf.paragraphs)
    if not paras:
        return ""
    anchor = {1: "center", 2: "flex-end"}.get(
        getattr(tf.vertical_anchor, "value", None), "flex-start")
    return ('<div class="shp txt" style="%s;justify-content:%s">%s</div>'
            % (pos, anchor, paras))

def notes_of(slide):
    try:
        if slide.has_notes_slide:
            return (slide.notes_slide.notes_text_frame.text or "").strip()
    except Exception:
        pass
    return ""

def convert(path, slug=None):
    prs = Presentation(path)
    W, H = prs.slide_width, prs.slide_height
    slug = slug or os.path.splitext(os.path.basename(path))[0]
    media_dir = os.path.join(OUT, "media", slug)
    if os.path.isdir(media_dir):
        shutil.rmtree(media_dir)
    os.makedirs(media_dir, exist_ok=True)
    media = (media_dir, "media/" + slug)
    counter = [0]
    frags, narr = {}, {}
    for i, slide in enumerate(prs.slides, start=1):
        shapes = [s for s in slide.shapes]
        parts = []
        for s_ in shapes:
            counter[0] += 1
            parts.append(shape_html(s_, W, H, media, counter[0]))
        body = "".join(parts)
        frags[i] = '<div class="slide">%s</div>' % body
        n = notes_of(slide)
        # '---' in the notes marks reveal steps
        parts = [p.strip() for p in re.split(r"^\s*---\s*$", n, flags=re.M) if p.strip()]
        narr[i] = "\n\n".join(parts) if parts else n
    title = os.path.splitext(os.path.basename(path))[0].replace("-", " ").replace("_", " ")
    return title, frags, narr, (W / H if H else 16 / 9)

# ---------------------------------------------------------------- html output
BASE_URL = os.environ.get(
    "SLIDES_BASE_URL",
    "https://aduyinuo.github.io/yinuo-d-is-an-open-book/slides")


def build_html(title, frags, narr, ratio, slug):
    tpl = open(os.path.join(HERE, "viewer_template.html"), encoding="utf-8").read()
    self_url = "%s/%s.html" % (BASE_URL, slug)
    return (tpl.replace("__TITLE__", html.escape(title, quote=False))
               .replace("__RATIO__", "%.4f" % ratio)
               .replace("__SELF__", self_url)
               .replace("__BASE__", BASE_URL)
               .replace("__SLUG__", slug)
               .replace("__FRAGS__", json.dumps(frags, ensure_ascii=False))
               .replace("__NARR__",  json.dumps(narr,  ensure_ascii=False)))


def write_oembed(slug, title, ratio):
    """Static oEmbed document so embedders render the deck as a live frame."""
    d = os.path.join(OUT, "oembed")
    os.makedirs(d, exist_ok=True)
    self_url = "%s/%s.html" % (BASE_URL, slug)
    w, h = 1280, 900
    doc = {
        "version": "1.0",
        "type": "rich",
        "provider_name": "Open Book Lab",
        "provider_url": BASE_URL,
        "title": title,
        "width": w,
        "height": h,
        "thumbnail_url": "%s/media/%s/cover.png" % (BASE_URL, slug),
        "thumbnail_width": 1280,
        "thumbnail_height": 720,
        "html": ('<iframe src="%s" width="%d" height="%d" frameborder="0" '
                 'allowfullscreen style="border:0;max-width:100%%"></iframe>'
                 % (self_url, w, h)),
    }
    with open(os.path.join(d, slug + ".json"), "w", encoding="utf-8", newline="\n") as fh:
        json.dump(doc, fh, ensure_ascii=False, indent=2)


def write_cover(slug, title):
    """A simple cover image so link previews are not blank."""
    try:
        from PIL import Image, ImageDraw, ImageFont
    except Exception:
        return
    W, H = 1280, 720
    im = Image.new("RGB", (W, H), "#22452f")
    d = ImageDraw.Draw(im)
    def font(px):
        for f in ("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",
                  "/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf",
                  "C:/Windows/Fonts/segoeuib.ttf"):
            if os.path.exists(f):
                return ImageFont.truetype(f, px)
        return ImageFont.load_default()
    f1 = font(64)
    words, lines, cur = title.split(), [], ""
    for w_ in words:
        t = (cur + " " + w_).strip()
        if d.textlength(t, font=f1) > W - 160:
            lines.append(cur); cur = w_
        else:
            cur = t
    lines.append(cur)
    y = H // 2 - (len(lines) * 76) // 2
    for ln in lines[:4]:
        d.text((80, y), ln, font=f1, fill="#ffffff"); y += 76
    d.rectangle([80, y + 24, 200, y + 32], fill="#93b294")
    md = os.path.join(OUT, "media", slug)
    os.makedirs(md, exist_ok=True)
    im.save(os.path.join(md, "cover.png"))

def main(argv):
    os.makedirs(OUT, exist_ok=True)
    os.makedirs(SRC, exist_ok=True)
    files = argv[1:] or [os.path.join(SRC, f) for f in sorted(os.listdir(SRC))
                         if f.lower().endswith((".pptx", ".potx"))]
    if not files:
        print("No .pptx found in slides-source/. Put a deck there and run again.")
        return 1
    made = []
    for f in files:
        if not os.path.isabs(f) and not os.path.exists(f):
            f = os.path.join(SRC, f)
        f = os.path.abspath(f)
        slug = os.path.splitext(os.path.basename(f))[0]
        title, frags, narr, ratio = convert(f, slug)
        out = os.path.join(OUT, slug + ".html")
        open(out, "w", encoding="utf-8", newline="\n").write(
            build_html(title, frags, narr, ratio, slug))
        write_oembed(slug, title, ratio)
        write_cover(slug, title)
        made.append((slug, len(frags), sum(1 for v in narr.values() if v)))
    for slug, n, withnotes in made:
        print("%-40s %2d slides, %d with narration" % (slug, n, withnotes))
        print("   embed: https://aduyinuo.github.io/yinuo-d-is-an-open-book/slides/%s.html" % slug)
    return 0

if __name__ == "__main__":
    sys.exit(main(sys.argv))
