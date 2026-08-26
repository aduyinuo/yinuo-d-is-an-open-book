#!/usr/bin/env python3
r"""
to-pptx.py --- turn a compiled beamer PDF into a PowerPoint file.

    python tools/to-pptx.py main.pdf
    python tools/to-pptx.py main.pdf -o talk.pptx --dpi 220 --notes-from main.tex

WHAT THIS DOES, AND WHAT IT DOES NOT

Each PDF page becomes one PowerPoint slide, with the page placed as a
full-bleed image. Slide size is taken from the PDF, so 4:3 and 16:9 both
come out right. Speaker notes can be carried across from the .tex source.

The text is NOT editable in PowerPoint. It cannot be: the slides are set
by LaTeX, with LaTeX's maths, fonts, spacing and TikZ drawings, and none
of that has a faithful PowerPoint equivalent. Anything claiming otherwise
is either re-typing your deck or mangling it. So this produces a deck you
can *present* from PowerPoint, hand to someone who requires .pptx, or
drop individual slides out of -- not one you can rewrite there.

If you need editable text, the honest options are to keep editing the
LaTeX, or to rebuild the deck natively in PowerPoint.

A page-level text layer is attached to each slide's notes (after any
speaker note), so the deck stays searchable and quotable.

ANIMATIONS

An overlay prints as several PDF pages, so an animated frame becomes
several near-identical slides. Usually you want the flattened version:

    pdflatex '\PassOptionsToClass{handout}{beamer}\input{main}'

or add "handout" to \documentclass and rebuild, then convert that PDF.

REQUIREMENTS

    pip install pymupdf python-pptx
"""

import argparse
import io
import os
import re
import sys

try:
    import pymupdf
except ImportError:
    sys.exit("Missing PyMuPDF.  pip install pymupdf")

try:
    from pptx import Presentation
    from pptx.util import Emu
except ImportError:
    sys.exit("Missing python-pptx.  pip install python-pptx")


EMU_PER_PT = 12700


def _balanced(src, open_index):
    """Return the text inside a brace group whose "{" is at open_index."""
    i = open_index + 1
    depth = 1
    while i < len(src) and depth:
        c = src[i]
        if c == "\\":                 # skip an escaped character
            i += 2
            continue
        if c == "{":
            depth += 1
        elif c == "}":
            depth -= 1
        i += 1
    return src[open_index + 1:i - 1], i


def _detex(text):
    text = re.sub(r"%.*", "", text)
    text = re.sub(r"\\[a-zA-Z]+\*?", "", text)
    text = text.replace("{", "").replace("}", "").replace("\\", "")
    return re.sub(r"\s+", " ", text).strip()


def extract_notes(tex_path):
    """Return [(frame_title, note_text)] in source order.

    Notes are tied to the title of the frame they sit in, not to a
    position in the file. Position alone is wrong the moment a frame
    spans several overlay pages, or a backup slide appears, or a frame
    carries no note -- the notes then slide onto the wrong slides, which
    is worse than having none.
    """
    try:
        src = io.open(tex_path, encoding="utf-8", errors="replace").read()
    except OSError as exc:
        print("  could not read %s: %s" % (tex_path, exc))
        return []

    # Split the source into frames, remembering each frame's title.
    frames = []
    for m in re.finditer(r"\\begin\{frame\}", src):
        k = m.end()
        # skip [options]
        while k < len(src) and src[k] in " \t\n":
            k += 1
        if k < len(src) and src[k] == "[":
            depth, k = 1, k + 1
            while k < len(src) and depth:
                if src[k] == "[":
                    depth += 1
                elif src[k] == "]":
                    depth -= 1
                k += 1
        while k < len(src) and src[k] in " \t\n":
            k += 1
        title = ""
        if k < len(src) and src[k] == "{":
            title, _ = _balanced(src, k)
        end = src.find(r"\end{frame}", m.end())
        if end == -1:
            end = len(src)
        frames.append((_detex(title), src[m.end():end]))

    notes = []
    for title, body in frames:
        found = []
        for nm in re.finditer(r"\\(?:speakernote|note)\s*(?:<[^>]*>)?\s*\{", body):
            text, _ = _balanced(body, nm.end() - 1)
            text = _detex(text)
            if text:
                found.append(text)
        if found:
            notes.append((title, "\n\n".join(found)))
    return notes


def convert(pdf_path, out_path, dpi, notes_from, quiet=False):
    doc = pymupdf.open(pdf_path)
    if doc.page_count == 0:
        sys.exit("%s has no pages." % pdf_path)

    notes = extract_notes(notes_from) if notes_from else []
    if notes_from and not quiet:
        print("  %d speaker note(s) found in %s" % (len(notes), notes_from))

    # Match each note to the FIRST page whose text carries its frame
    # title. Overlay pages repeat the title, so the note lands on the
    # page the frame starts on.
    page_note = {}
    unmatched = []
    if notes:
        page_text = [pg.get_text() for pg in doc]
        for title, body in notes:
            if not title:
                unmatched.append(title)
                continue
            needle = " ".join(title.split())
            hit = None
            for n, txt in enumerate(page_text):
                if needle and needle in " ".join(txt.split()):
                    hit = n
                    break
            if hit is None:
                unmatched.append(title)
            else:
                page_note.setdefault(hit, []).append(body)

    first = doc[0]
    w_pt, h_pt = first.rect.width, first.rect.height

    prs = Presentation()
    prs.slide_width = Emu(int(round(w_pt * EMU_PER_PT)))
    prs.slide_height = Emu(int(round(h_pt * EMU_PER_PT)))
    blank = prs.slide_layouts[6]          # 6 is the blank layout

    zoom = dpi / 72.0
    matrix = pymupdf.Matrix(zoom, zoom)

    for index, page in enumerate(doc):
        if (page.rect.width, page.rect.height) != (w_pt, h_pt):
            print("  note: page %d has a different size; it will be "
                  "stretched to match page 1" % (index + 1))

        pix = page.get_pixmap(matrix=matrix, alpha=False)
        image = io.BytesIO(pix.tobytes("png"))

        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(image, 0, 0,
                                 width=prs.slide_width,
                                 height=prs.slide_height)

        parts = []
        if index in page_note:
            parts.extend(page_note[index])
        text = page.get_text().strip()
        if text:
            parts.append("--- slide text ---\n" + text)
        if parts:
            slide.notes_slide.notes_text_frame.text = "\n\n".join(parts)

        if not quiet and (index + 1) % 10 == 0:
            print("  %d/%d pages" % (index + 1, doc.page_count))

    n_pages = doc.page_count
    prs.save(out_path)
    doc.close()

    size_mb = os.path.getsize(out_path) / (1024 * 1024)
    if not quiet:
        print("wrote %s" % out_path)
        print("  %d slides, %.0f x %.0f pt, %d dpi, %.1f MB"
              % (len(prs.slides), w_pt, h_pt, dpi, size_mb))
        if notes:
            print("  %d of %d note(s) placed by frame title"
                  % (len(notes) - len(unmatched), len(notes)))
            for t in unmatched:
                print("    unplaced (title not found on any page): %r"
                      % (t or "<untitled frame>"))


def main():
    ap = argparse.ArgumentParser(
        description="Turn a compiled beamer PDF into a PowerPoint file.")
    ap.add_argument("pdf", help="the compiled PDF, e.g. main.pdf")
    ap.add_argument("-o", "--out", help="output .pptx (default: alongside the PDF)")
    ap.add_argument("--dpi", type=int, default=200,
                    help="render resolution, default 200. Use 300 for print, "
                         "150 to keep the file small.")
    ap.add_argument("--notes-from", metavar="TEX",
                    help="a .tex file to lift \\speakernote{} / \\note{} from")
    ap.add_argument("-q", "--quiet", action="store_true")
    args = ap.parse_args()

    if not os.path.exists(args.pdf):
        sys.exit("No such file: %s" % args.pdf)
    if args.dpi < 50 or args.dpi > 600:
        sys.exit("--dpi should be between 50 and 600.")

    out = args.out or os.path.splitext(args.pdf)[0] + ".pptx"
    convert(args.pdf, out, args.dpi, args.notes_from, args.quiet)


if __name__ == "__main__":
    main()
