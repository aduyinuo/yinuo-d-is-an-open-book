#!/usr/bin/env python3
r"""
from-pptx.py --- rebuild a PowerPoint deck as a Lucid beamer deck.

    python tools/from-pptx.py talk.pptx -o rebuilt/

WHAT IT RECOVERS

    * the title of every slide, and its body text as bullets
    * every speaker note
    * every embedded image, at original resolution, placed on its slide
    * any table, as a lucidtable
    * the aspect ratio

WHAT IT CANNOT RECOVER

Diagrams that PowerPoint drew from its own shapes -- boxes, arrows,
connectors, grouped constructions -- are not pictures. They live as
drawing instructions with no image behind them, and there is nothing to
extract. Those slides come out with their text and a marked placeholder
saying what was on them, so you can redraw them in TikZ or paste a
screenshot.

This is the same limit that mangles a deck when it is pushed through
pdf-to-HTML converters: they keep the embedded pictures, drop the drawn
shapes, and lose the layout. The difference is that this script tells
you which slides lost something instead of quietly shipping them.

Run it, read the report at the end, then fix the listed slides by hand.

REQUIREMENTS

    pip install python-pptx
"""

import argparse
import os
import re
import sys

try:
    from pptx import Presentation
    from pptx.util import Emu
except ImportError:
    sys.exit("Missing python-pptx.  pip install python-pptx")


# LaTeX special characters, in an order that does not double-escape.
_TEX = [
    ("\\", r"\textbackslash{}"),
    ("&", r"\&"), ("%", r"\%"), ("$", r"\$"), ("#", r"\#"),
    ("_", r"\_"), ("{", r"\{"), ("}", r"\}"),
    ("~", r"\textasciitilde{}"), ("^", r"\textasciicircum{}"),
]


# Characters that actually turn up in slide decks, mapped to something
# pdflatex can set. Everything unmapped is dropped, not turned into "?",
# because a stray "?" reads as a real question mark in the text.
_UNI = [
    ("\u2019", "'"), ("\u2018", "`"), ("\u201c", "``"), ("\u201d", "''"),
    ("\u2013", "--"), ("\u2014", "---"), ("\u2026", r"\ldots{}"),
    ("\u00a0", " "), ("\u2192", r"$\rightarrow$"), ("\u2190", r"$\leftarrow$"),
    ("\u21d2", r"$\Rightarrow$"), ("\u00d7", r"$\times$"),
    ("\u2265", r"$\geq$"), ("\u2264", r"$\leq$"), ("\u2260", r"$\neq$"),
    ("\u2248", r"$\approx$"), ("\u00b1", r"$\pm$"), ("\u221e", r"$\infty$"),
    ("\u03b1", r"$\alpha$"), ("\u03b2", r"$\beta$"), ("\u03b3", r"$\gamma$"),
    ("\u03bb", r"$\lambda$"), ("\u03bc", r"$\mu$"), ("\u03c0", r"$\pi$"),
    ("\u03c3", r"$\sigma$"), ("\u03b8", r"$\theta$"), ("\u03b5", r"$\epsilon$"),
    ("\u2022", ""), ("\u25cf", ""), ("\u25aa", ""), ("\u2013", "--"),
    ("\u00b0", r"$^\circ$"), ("\u00ae", ""), ("\u2122", ""),
    # Non-ASCII hyphens and dashes. Dropping these silently welds words
    # together -- a non-breaking hyphen turns "Cyber-Deception" into
    # "CyberDeception" -- so every one of them is mapped.
    ("\u2010", "-"), ("\u2011", "-"), ("\u2012", "--"), ("\u2015", "---"),
    ("\u2212", "-"), ("\u00ad", "-"), ("\u2044", "/"),
]


def tex(s):
    """Escape a plain string for LaTeX."""
    if not s:
        return ""
    for a, b in _TEX:
        s = s.replace(a, b)
    # Strip decorative leading glyphs BEFORE mapping. Doing it after
    # would eat the leading dollar of an arrow substitution and leave
    # unbalanced math behind.
    s = re.sub(r"^[^\w]{1,4}\s*(?=[A-Za-z0-9])", "", s)
    for a, b in _UNI:
        s = s.replace(a, b)
    # Anything still outside ASCII would need inputenc to survive
    # pdflatex, so drop it rather than break the build.
    return "".join(c if ord(c) < 128 else "" for c in s).strip()


def is_pic(shape):
    return str(shape.shape_type).startswith("PICTURE")


def is_drawn(shape):
    kind = str(shape.shape_type)
    return (kind.startswith("AUTO_SHAPE") or kind.startswith("GROUP")
            or kind.startswith("LINE") or kind.startswith("FREEFORM")
            or kind.startswith("CONNECTOR"))


def slide_title(slide):
    if slide.shapes.title is not None and slide.shapes.title.has_text_frame:
        t = slide.shapes.title.text_frame.text.strip()
        if t:
            return " ".join(t.split())
    return ""


def body_bullets(slide, skip_shape, limit=6):
    """Text from every non-title text frame, as bullet strings."""
    out = []
    for sh in slide.shapes:
        if sh is skip_shape or not sh.has_text_frame:
            continue
        for para in sh.text_frame.paragraphs:
            line = " ".join(r.text for r in para.runs).strip()
            line = " ".join(line.split())
            if len(line) > 1:
                out.append(line)
    # drop duplicates, keep order
    seen, uniq = set(), []
    for line in out:
        key = line.lower()
        if key not in seen:
            seen.add(key)
            uniq.append(line)
    return uniq[:limit]


def save_images(slide, index, img_dir, rel_dir):
    """Write out the slide's pictures; return (relpaths, widths_fraction)."""
    paths = []
    for n, sh in enumerate(s for s in slide.shapes if is_pic(s)):
        try:
            image = sh.image
        except Exception:
            continue
        ext = image.ext or "png"
        name = "s%02d_%d.%s" % (index, n, ext)
        with open(os.path.join(img_dir, name), "wb") as fh:
            fh.write(image.blob)
        paths.append("%s/%s" % (rel_dir, name))
    return paths


def slide_tables(slide):
    out = []
    for sh in slide.shapes:
        if getattr(sh, "has_table", False) and sh.has_table:
            rows = []
            for r in sh.table.rows:
                rows.append([" ".join(c.text.split()) for c in r.cells])
            if rows:
                out.append(rows)
    return out


def render_table(rows):
    ncol = max(len(r) for r in rows)
    rows = [r + [""] * (ncol - len(r)) for r in rows]
    spec = "l" * ncol
    lines = [r"  \begin{lucidtable}", r"  \begin{tabular}{%s}" % spec,
             r"    \toprule"]
    head = " & ".join(tex(c) for c in rows[0])
    lines.append(r"    \thead{%s} \\" % head)
    lines.append(r"    \midrule")
    for r in rows[1:]:
        lines.append("    " + " & ".join(tex(c) for c in r) + r" \\")
    lines += [r"    \bottomrule", r"  \end{tabular}", r"  \end{lucidtable}"]
    return "\n".join(lines)


# ---------------------------------------------------------------------
#  Redrawing PowerPoint's own shapes
#
#  A drawn shape carries no image, but it does carry geometry, a shape
#  kind, a colour and its text. That is enough to lay the diagram out
#  again in TikZ: boxes land where they were, arrows still point the
#  same way, labels keep their words. It is a reconstruction, not a
#  photograph -- curves, shadows and gradients are not reproduced -- but
#  it beats a placeholder saying something used to be here.
# ---------------------------------------------------------------------

EMU_PER_PT = 12700.0


def _pt(v):
    return (v or 0) / EMU_PER_PT


def _rgb(shape):
    """Fill colour as a LaTeX-safe hex string, or None for theme default."""
    try:
        f = shape.fill
        if f.type is not None and int(f.type) == 1:
            return str(f.fore_color.rgb)
    except Exception:
        pass
    return None


def _line_rgb(shape):
    try:
        c = shape.line.color
        if c and c.type is not None:
            return str(c.rgb)
    except Exception:
        pass
    return None


_A_NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"


def _flip(shape):
    """(flipH, flipV) for a shape.

    PowerPoint stores a connector as a bounding box plus two flip flags,
    not as a start and an end point. Ignore the flags and every arrow
    gets drawn along the same diagonal, so roughly half of them come out
    pointing the wrong way -- which is the single most visible way a
    rebuilt diagram can be wrong.
    """
    try:
        x = shape.element.find(".//" + _A_NS + "xfrm")
        if x is None:
            return False, False
        return x.get("flipH") == "1", x.get("flipV") == "1"
    except Exception:
        return False, False


def _flatten(shapes, dx=0.0, dy=0.0):
    """Walk groups so nested shapes come back with absolute positions."""
    out = []
    for sh in shapes:
        kind = str(sh.shape_type)
        if kind.startswith("GROUP"):
            try:
                out.extend(_flatten(sh.shapes, dx + _pt(sh.left), dy + _pt(sh.top)))
            except Exception:
                pass
            continue
        out.append((sh, dx, dy))
    return out


def _drawn_items(slide):
    """Shapes that make up a drawing, with absolute positions."""
    out = []
    for sh, dx, dy in _flatten(slide.shapes):
        kind = str(sh.shape_type)
        if kind.startswith("PICTURE") or kind.startswith("TABLE"):
            continue
        text = ""
        if sh.has_text_frame:
            text = " ".join(sh.text_frame.text.split())
        is_box = (kind.startswith("AUTO_SHAPE") or kind.startswith("FREEFORM"))
        is_line = (kind.startswith("LINE") or kind.startswith("CONNECTOR"))
        is_label = kind.startswith("TEXT_BOX") and text
        if not (is_box or is_line or is_label):
            continue
        w, h = _pt(sh.width), _pt(sh.height)
        if w <= 0 and h <= 0:
            continue
        out.append({
            "sh": sh, "text": text, "box": is_box, "line": is_line,
            "x": _pt(sh.left) + dx, "y": _pt(sh.top) + dy, "w": w, "h": h,
        })
    return out


def diagram_texts(slide):
    """Every string that will appear inside the rebuilt diagram."""
    return set(i["text"] for i in _drawn_items(slide) if i["text"])


def shapes_to_tikz(slide, escape, target_w=290.0, target_h=125.0):
    """Rebuild a slide's drawn shapes as TikZ, sized to fit the slide.

    Coordinates are scaled in Python and emitted as absolute points, so
    the picture fits without a tikz `scale=`. That matters: `scale` with
    `transform shape` shrinks the labels too, and a diagram whose text
    has been scaled to two points is no better than the placeholder it
    replaced.
    """
    items = _drawn_items(slide)
    if not items:
        return ""

    x0 = min(i["x"] for i in items)
    y0 = min(i["y"] for i in items)
    x1 = max(i["x"] + i["w"] for i in items)
    y1 = max(i["y"] + i["h"] for i in items)
    bw, bh = max(x1 - x0, 1.0), max(y1 - y0, 1.0)
    k = min(target_w / bw, target_h / bh, 1.0)

    def X(v):
        return (v - x0) * k

    def Y(v):
        return (y1 - v) * k          # flip: PowerPoint y grows downward

    palette = {}

    def cname(hexrgb, fallback):
        if not hexrgb:
            return fallback
        h = str(hexrgb).upper()[-6:]
        if len(h) != 6 or any(c not in "0123456789ABCDEF" for c in h):
            return fallback
        palette.setdefault(h, "ppt%s" % h)
        return palette[h]

    body = []
    for i in items:
        sh = i["sh"]
        ax, ay = X(i["x"]), Y(i["y"])
        bx, by = X(i["x"] + i["w"]), Y(i["y"] + i["h"])
        cx, cy = (ax + bx) / 2.0, (ay + by) / 2.0
        label = escape(i["text"])

        if i["line"]:
            col = cname(_line_rgb(sh), "lucidMuted")
            fh, fv = _flip(sh)
            sx, ex = (bx, ax) if fh else (ax, bx)
            sy, ey = (by, ay) if fv else (ay, by)
            body.append(r"    \draw[->,%s,line width=0.5pt] (%.1f,%.1f) -- (%.1f,%.1f);"
                        % (col, sx, sy, ex, ey))
        elif i["box"]:
            fc = cname(_rgb(sh), "lucidBand!50")
            dc = cname(_line_rgb(sh), "lucidAccent")
            body.append(r"    \draw[fill=%s,draw=%s,line width=0.4pt,rounded corners=2pt]"
                        r" (%.1f,%.1f) rectangle (%.1f,%.1f);" % (fc, dc, ax, ay, bx, by))
            if label:
                body.append(r"    \node[text width=%.1fpt,align=center,font=\tiny,"
                            r"text=lucidInk] at (%.1f,%.1f) {%s};"
                            % (max((bx - ax) - 4, 14), cx, cy, label))
        elif label:
            body.append(r"    \node[text width=%.1fpt,align=left,anchor=north west,"
                        r"font=\tiny,text=lucidInk] at (%.1f,%.1f) {%s};"
                        % (max(bx - ax, 24), ax, ay, label))

    defs = "".join("  \\definecolor{%s}{HTML}{%s}\n" % (n, h)
                   for h, n in sorted(palette.items()))
    return (defs
            + "  \\begin{center}\n"
            + "  \\begin{tikzpicture}[x=1pt,y=1pt]\n"
            + "\n".join(body) + "\n"
            + "  \\end{tikzpicture}\n  \\end{center}")


def build(pptx_path, out_dir, max_bullets):
    prs = Presentation(pptx_path)
    ratio = prs.slide_width / float(prs.slide_height)
    aspect = "169" if ratio > 1.6 else "43"

    img_dir = os.path.join(out_dir, "figures")
    os.makedirs(img_dir, exist_ok=True)

    page_w = prs.slide_width / EMU_PER_PT
    page_h = prs.slide_height / EMU_PER_PT
    body = []
    needs_work = []
    rebuilt = []
    lost_chars = 0

    for i, slide in enumerate(prs.slides, 1):
        title = slide_title(slide)
        title_shape = slide.shapes.title
        bullets = body_bullets(slide, title_shape, max_bullets)
        # Text that the rebuilt diagram will draw must not also be
        # listed as a bullet, or every label appears twice.
        if not [x for x in slide.shapes if is_pic(x)]:
            inart = set(' '.join(t.split()).lower()
                        for t in diagram_texts(slide))
            if inart:
                # A box label often arrives split across paragraphs, so
                # "Instructions," and "Quiz, Practice" are separate
                # bullets while the box says "Instructions, Quiz,
                # Practice". Exact matching misses those; containment
                # catches them.
                joined = ' | '.join(inart)
                bullets = [b for b in bullets
                           if ' '.join(b.split()).lower() not in inart
                           and ' '.join(b.split()).lower() not in joined]
        images = save_images(slide, i, img_dir, "figures")
        tables = slide_tables(slide)
        drawn = sum(1 for sh in slide.shapes if is_drawn(sh))

        note = ""
        if slide.has_notes_slide:
            note = " ".join(slide.notes_slide.notes_text_frame.text.split())

        # A slide whose picture is really the whole slide gets the figure
        # its own frame; otherwise text leads and the figure follows.
        head = title or (bullets[0] if bullets else "Slide %d" % i)
        if not title and bullets:
            bullets = bullets[1:]
        # A slide's title is very often repeated in a separate text box,
        # so it arrives as a bullet too. Printing it twice is the most
        # visible artefact of a naive conversion.
        key = " ".join(head.lower().split())
        bullets = [b for b in bullets
                   if " ".join(b.lower().split()) != key]

        # A recovered slide is often taller than one beamer frame:
        # PowerPoint text boxes hold more than this type size does. Let
        # dense frames continue onto a second page rather than overprint
        # the footer. Truncating instead would throw away text the
        # author wrote, which is worse than an extra page.
        dense = len(bullets) > 3 or tables or (bullets and images)
        opt = "[allowframebreaks]" if dense else ""

        lines = []
        lines.append("%% ---- slide %d " % i + "-" * 50)
        lines.append(r"\begin{frame}%s{%s}" % (opt, tex(head)))

        if bullets:
            lines.append(r"  \begin{itemize}")
            for b in bullets:
                lines.append(r"    \item %s" % tex(b))
            lines.append(r"  \end{itemize}")

        for t in tables:
            lines.append(render_table(t))

        if images:
            if len(images) == 1:
                w = "0.62\\textwidth" if bullets else "0.82\\textwidth"
                lines.append(r"  \slidefigure[%s]{%s}" % (w, images[0]))
            else:
                each = min(0.42, 0.9 / min(len(images), 3))
                lines.append(r"  \begin{center}")
                for pth in images[:3]:
                    lines.append(r"    \includegraphics[width=%.2f\textwidth,"
                                 r"height=0.34\textheight,keepaspectratio]{%s}\hfill"
                                 % (each, pth))
                lines.append(r"  \end{center}")
                if len(images) > 3:
                    lines.append(r"  %% %d more image(s) in figures/, not placed"
                                 % (len(images) - 3))

        if drawn and not images:
            art = shapes_to_tikz(slide, tex)
            if art:
                lines.append(art)
                rebuilt.append((i, head, drawn))
            else:
                lines.append(r"  \vfill")
                lines.append(r"  \begin{center}")
                lines.append(r"    \fbox{\begin{minipage}{0.8\textwidth}\centering")
                lines.append(r"      \color{lucidMuted}\small This slide carried a "
                             r"drawing with no recoverable geometry.\end{minipage}}")
                lines.append(r"  \end{center}")
                lines.append(r"  \vfill")
                needs_work.append((i, head, drawn))

        if note:
            lines.append(r"  \speakernote{%s}" % tex(note))

        lines.append(r"\end{frame}")
        lines.append("")
        body.append("\n".join(lines))

    return aspect, body, needs_work, len(prs.slides), rebuilt


PREAMBLE = r"""%% =====================================================================
%%  %(name)s --- rebuilt from %(src)s with tools/from-pptx.py
%%
%%  Text, speaker notes, tables and embedded images were recovered from
%%  the PowerPoint file. Diagrams that PowerPoint drew from its own
%%  shapes could not be: they are drawing instructions, not pictures.
%%  Slides that lost one carry a REDRAW box -- see the list at the end
%%  of this file.
%%
%%  Build:  pdflatex %(name)s   (twice)
%% =====================================================================

\documentclass[11pt,aspectratio=%(aspect)s]{beamer}
\usetheme{lucid}

\usepackage{graphicx}

\title{%(title)s}
\date{}

%% \author{Your Name}
%% \institute{Your Institution}

\begin{document}

\begin{frame}[plain]
  \titlepage
\end{frame}

"""


def main():
    ap = argparse.ArgumentParser(
        description="Rebuild a .pptx as a Lucid beamer deck.")
    ap.add_argument("pptx")
    ap.add_argument("-o", "--out", default="rebuilt",
                    help="output directory (default: rebuilt/)")
    ap.add_argument("--title", default=None, help="deck title")
    ap.add_argument("--max-bullets", type=int, default=6)
    args = ap.parse_args()

    if not os.path.exists(args.pptx):
        sys.exit("No such file: %s" % args.pptx)
    os.makedirs(args.out, exist_ok=True)

    name = re.sub(r"[^a-zA-Z0-9]+", "-",
                  os.path.splitext(os.path.basename(args.pptx))[0]).strip("-")
    aspect, body, needs_work, total, rebuilt = build(
        args.pptx, args.out, args.max_bullets)

    title = args.title or name.replace("-", " ").title()
    text = PREAMBLE % {"name": name, "src": os.path.basename(args.pptx),
                       "aspect": aspect, "title": tex(title)}
    text += "\n".join(body)

    if needs_work:
        text += "\n%% --- slides whose diagram must be redrawn by hand ---\n"
        for i, head, n in needs_work:
            text += "%%   slide %-3d %-52s (%d shapes)\n" % (i, head[:52], n)

    text += "\n\\end{document}\n"

    out_tex = os.path.join(args.out, name + ".tex")
    with open(out_tex, "w", encoding="utf-8", newline="\n") as fh:
        fh.write(text)

    print("wrote %s" % out_tex)
    print("  %d slides, aspect %s" % (total, aspect))
    print("  images extracted to %s/figures/" % args.out)
    print("  %d diagram(s) rebuilt as TikZ from the original shapes"
          % len(rebuilt))
    print("  %d slide(s) still need a diagram by hand:" % len(needs_work))
    for i, head, n in needs_work[:15]:
        safe = head[:50].encode("ascii", "replace").decode("ascii")
        print("    slide %-3d %-50s (%d shapes)" % (i, safe, n))
    if len(needs_work) > 15:
        print("    ... and %d more, all listed at the end of the .tex"
              % (len(needs_work) - 15))
    print("\nCopy beamerthemelucid.sty into %s/ and run pdflatex twice."
          % args.out)


if __name__ == "__main__":
    main()
